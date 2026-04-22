import io
import math
from docx import Document
import pdfplumber

def extract_text_from_pdf(file_bytes: bytes) -> str:
    text_parts = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
    print(text_parts)
    return "\n".join(text_parts)

def extract_text_from_docx(file_bytes: bytes) -> str:
    """
    Extracts all available text from a DOCX file, including paragraphs and table cells.
    
    Args:
        file_bytes (bytes): The raw binary content of a DOCX file.

    Returns:
        str: All extracted text, separated by newlines.
    """
    document = Document(io.BytesIO(file_bytes))
    text_parts = []

    # Extract all paragraph text
    for para in document.paragraphs:
        if para.text.strip():
            text_parts.append(para.text.strip())

    # Extract text from tables
    for table in document.tables:
        for row in table.rows:
            for cell in row.cells:
                cell_text = cell.text.strip()
                if cell_text:
                    text_parts.append(cell_text)

    print(text_parts)
    return "\n".join(text_parts)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """
    Extracts all text from a PowerPoint PPTX file.
    
    Args:
        file_bytes (bytes): The raw binary content of a PPTX file.

    Returns:
        str: All extracted text, separated by newlines.
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ValueError("python-pptx is required to process PowerPoint files. Install with: pip install python-pptx")
    
    prs = Presentation(io.BytesIO(file_bytes))
    text_parts = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())
            # Handle tables in slides
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_parts.append(cell.text.strip())
    
    print(text_parts)
    return "\n".join(text_parts)


def extract_text(file_bytes: bytes, name: str) -> str:
    ext = name.lower().split('.')[-1]
    if ext == "pdf":
        return extract_text_from_pdf(file_bytes)
    elif ext == "docx":
        return extract_text_from_docx(file_bytes)
    elif ext in ("ppt", "pptx"):
        return extract_text_from_pptx(file_bytes)
    else:
        raise ValueError(f"Unsupported file extension: {ext}")


def extract_question_type_mix(config: dict, manual_questions: list) -> dict:
    allowed_keys = ["mcq", "coding", "architecture", "scenario", "screening"]

    # Step 1: base from config
    result = {}
    config = config or {"mcq": 6, "architecture": 2, "coding": 2, "scenario": 0, "screening": 0}

    for key in allowed_keys:
        try:
            value = int(config.get(key, 0))
        except (ValueError, TypeError):
            value = 0

        result[key] = value

    # Step 2: add manual questions
    for q in manual_questions or []:
        q_type = None

        # handle dict or ORM
        if isinstance(q, dict):
            q_type = q.get("type") or q.get("question_type")
        else:
            q_type = getattr(q, "question_type", None)

        if q_type in allowed_keys:
            result[q_type] += 1

    return result

def calculate_duration_minutes(config: dict) -> int:
    if not config :
        return 30  # default duration for unknown config
    
    time_map = {
        "mcq": 30,
        "architecture": 120,
        "coding": 600,
        "scenario": 120,
        "screening": 120,
    }

    total_seconds = 0

    #  1. From questionnaire_config
    for key in ["mcq", "coding", "architecture", "scenario", "screening"]:
        count = config.get(key, 0)

        try:
            count = int(count)
        except:
            count = 0

        total_seconds += count * time_map[key]
    
    if total_seconds == 0:
        return 30  # default duration for zero questions

    return math.ceil(total_seconds / 60)
